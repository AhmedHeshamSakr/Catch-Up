#!/usr/bin/env python3
"""Generate the Catch-Up product deck (Google-style, 10 slides, low-text/visual).

Run:  uv run --with python-pptx python presentation/build_deck.py
Out:  presentation/Catch-Up.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
LOGOS = os.path.join(HERE, "assets", "logos")
ICONS = os.path.join(HERE, "assets", "icons")
OUT = os.path.join(HERE, "Catch-Up.pptx")

# ---- Google palette ---------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x20, 0x21, 0x24)   # grey 900 — headings
BODY  = RGBColor(0x3C, 0x40, 0x43)   # grey 800 — body
MUTE  = RGBColor(0x5F, 0x63, 0x68)   # grey 700 — captions
FAINT = RGBColor(0x80, 0x86, 0x8B)   # grey 600
LINE  = RGBColor(0xE8, 0xEA, 0xED)   # grey 200 — borders
TILE  = RGBColor(0xF8, 0xF9, 0xFA)   # grey 50  — soft fill
BLUE   = RGBColor(0x42, 0x85, 0xF4)
RED    = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xF9, 0xAB, 0x00)  # bright — fills/icons only
AMBER  = RGBColor(0xE3, 0x74, 0x00)  # readable warm accent for text/rules
GREEN  = RGBColor(0x34, 0xA8, 0x53)
FONT = "Arial"

EMU_IN = 914400


def icon(name):  # -> path
    return os.path.join(ICONS, f"{name}.png")


def logo(name):
    return os.path.join(LOGOS, f"{name}.png")


# ---- low-level helpers ------------------------------------------------------
def set_spc(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def text(slide, x, y, w, h, lines, *, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None, space_after=0,
         italic=False, font=FONT):
    """lines: str or list[str]; each becomes a paragraph."""
    if isinstance(lines, str):
        lines = [lines]
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = font
        f.color.rgb = color
        if spacing:
            set_spc(r, spacing)
    return tb


def bullets(slide, x, y, w, h, items, *, size, dash_color, text_color=BODY,
            gap=11, line_spacing=1.05, dash="—  ", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        rd = p.add_run()
        rd.text = dash
        rd.font.name = FONT
        rd.font.size = Pt(size)
        rd.font.bold = True
        rd.font.color.rgb = dash_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = text_color
    return tb


def flatten(shp):
    """Remove themed style (drop-shadow/fill/line) and force no effects -> flat."""
    from pptx.oxml.ns import qn
    sp = shp._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)
    spPr = shp._element.spPr
    for ex in spPr.findall(qn("a:effectLst")):
        spPr.remove(ex)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def rect(slide, x, y, w, h, fill, *, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        In(x), In(y), In(w), In(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.text_frame.paragraphs[0].text = ""
    flatten(shp)
    return shp


def oval(slide, x, y, d, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    flatten(shp)
    return shp


def pic(slide, path, x, y, *, h):
    """square assets: width auto-resolves equal to h."""
    return slide.shapes.add_picture(path, In(x), In(y), height=In(h))


def pic_w(shape):
    return shape.width / EMU_IN


def dots(slide, x, y, d=0.12, gap=0.08):
    for i, c in enumerate((BLUE, RED, YELLOW, GREEN)):
        oval(slide, x + i * (d + gap), y, d, c)


def page_no(slide, n=None):
    # auto: numbered by creation order, so inserting slides never desyncs
    num = len(prs.slides._sldIdLst)
    text(slide, 12.0, 6.96, 0.9, 0.3, f"{num:02d}", size=10, color=FAINT,
         align=PP_ALIGN.RIGHT, spacing=1.0)


def brandmark(slide):
    # tiny running header: green square + wordmark
    rect(slide, 0.9, 0.5, 0.13, 0.13, GREEN, rounded=True, radius=0.25)
    text(slide, 1.12, 0.45, 3.0, 0.3, "Catch-Up", size=12, color=INK, bold=True)


# ---- slide scaffolds --------------------------------------------------------
SW, SH = 13.333, 7.5
ML = 0.9


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def header(slide, kicker, kicker_color, headline, *, hsize=37, hwidth=8.2):
    brandmark(slide)
    text(slide, ML, 1.05, hwidth, 0.32, kicker.upper(), size=12.5,
         color=kicker_color, bold=True, spacing=2.2)
    rect(slide, ML + 0.02, 1.44, 0.5, 0.058, kicker_color)  # section-marker rule
    text(slide, ML, 1.70, hwidth, 1.6, headline, size=hsize, color=INK,
         bold=True, line_spacing=1.02)


def hero_tile(slide, icon_name, *, accent):
    tw = 2.95
    tx = SW - ML - tw
    ty = 2.35
    rect(slide, tx, ty, tw, tw, TILE, line=LINE, line_w=1.0, rounded=True, radius=0.07)
    rect(slide, tx, ty, 0.10, tw, accent, rounded=True, radius=0.5)
    ih = 1.5
    pic(slide, icon(icon_name), tx + (tw - ih) / 2 + 0.05, ty + (tw - ih) / 2, h=ih)


def hero_bullets(slide, items, *, size, color, gap=12):
    # vertically centered against the hero tile (y 2.35 .. 5.30)
    bullets(slide, ML, 2.35, 7.9, 2.95, items, size=size, dash_color=color,
            gap=gap, anchor=MSO_ANCHOR.MIDDLE)


def why_note(slide, accent, why, *, y=5.66):
    pw, ph = 0.66, 0.32
    rect(slide, ML, y, pw, ph, accent, rounded=True, radius=0.5)
    text(slide, ML, y - 0.005, pw, ph, "WHY", size=10.5, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.3)
    text(slide, ML + pw + 0.24, y - 0.03, 9.0, ph + 0.12, why, size=15, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)


def tnode(slide, rail_x, y, num, name, role, color, *, name_x=1.58, role_x=3.7):
    """One node on the vertical agent-tree timeline."""
    d = 0.30
    oval(slide, rail_x - d / 2, y - d / 2, d, color)
    text(slide, rail_x - d / 2, y - d / 2 - 0.012, d, d, str(num), size=12,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, name_x, y - 0.16, role_x - name_x - 0.05, 0.32, name, size=14,
         color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, role_x, y - 0.16, 8.1 - role_x, 0.32, role, size=12.5, color=MUTE,
         anchor=MSO_ANCHOR.MIDDLE)


def kv_row(slide, x, y, w, name, role, name_color):
    """A bold agent name + em-dash + muted role, on one line."""
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    for txt, sz, col, bold in ((name, 13.5, name_color, True),
                               ("  —  ", 13.5, FAINT, False),
                               (role, 13, BODY, False)):
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


# ---- build ------------------------------------------------------------------
prs = Presentation()
prs.slide_width = In(SW)
prs.slide_height = In(SH)

TECH = [
    (logo("adk"), "Google ADK"),
    (logo("gemini"), "Gemini"),
    (logo("fastapi"), "FastAPI"),
    (logo("nextjs"), "Next.js"),
    (logo("sqlite"), "SQLite"),
    (logo("firebase"), "Firestore"),
]


def logo_strip(slide, items, top, *, icon_h, cell_w, center_x, labels=False,
               label_color=MUTE, label_size=11):
    n = len(items)
    total = n * cell_w
    start = center_x - total / 2
    for i, (path, label) in enumerate(items):
        cell_left = start + i * cell_w
        sh = pic(slide, path, cell_left, top, h=icon_h)
        iw = pic_w(sh)
        sh.left = In(cell_left + (cell_w - iw) / 2)
        if labels:
            text(slide, cell_left, top + icon_h + 0.10, cell_w, 0.3, label,
                 size=label_size, color=label_color, align=PP_ALIGN.CENTER, bold=True)


# === Slide 1 — Title =========================================================
s = new_slide(prs)
dots(s, ML, 0.85)
text(s, ML, 2.45, 9, 0.35, "NEWS INTELLIGENCE AGENT", size=14, color=BLUE,
     bold=True, spacing=3.0)
# wordmark with green accent on "Up"
tb = s.shapes.add_textbox(In(ML), In(2.78), In(11), In(1.5))
tf = tb.text_frame
tf.word_wrap = True
for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
    setattr(tf, m, 0)
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "Catch-"
r1.font.name = FONT; r1.font.size = Pt(90); r1.font.bold = True; r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = "Up"
r2.font.name = FONT; r2.font.size = Pt(90); r2.font.bold = True; r2.font.color.rgb = GREEN
text(s, ML, 4.42, 10.5, 0.7,
     "Your world, caught up — collected, ranked, and summarized by AI.",
     size=21, color=MUTE)
# tech strip
text(s, ML, 5.55, 6, 0.3, "BUILT WITH", size=11, color=FAINT, bold=True, spacing=2.6)
logo_strip(s, TECH, 5.92, icon_h=0.5, cell_w=1.0, center_x=ML + 3.0, labels=False)
# footer
rect(s, ML, 6.78, SW - 2 * ML, 0.013, LINE)
text(s, ML, 6.92, 7, 0.3, "Ahmed Hesham  ·  AI Engineer", size=12, color=MUTE)
text(s, SW - ML - 5, 6.92, 5, 0.3, "Built on Google ADK · Gemini", size=12,
     color=MUTE, align=PP_ALIGN.RIGHT)

# === Slide 2 — Pain ==========================================================
s = new_slide(prs)
header(s, "The problem", RED, ["Too much news.", "Too little time."])
hero_bullets(s, [
    "News is scattered across feeds, sites, APIs, and video",
    "No fast catch-up in both Arabic and English",
    "Manual monitoring that just doesn't scale",
], size=18, color=RED, gap=16)
hero_tile(s, "hero_pain", accent=RED)
page_no(s, 2)

# === Slide 3 — Solution ======================================================
s = new_slide(prs)
header(s, "The solution", GREEN, ["One agent that", "catches you up."])
hero_bullets(s, [
    "Collects from everywhere",
    "Ranks what actually matters",
    "Summarizes in English + Arabic",
    "Delivers a digest — scheduled or on-demand",
], size=18, color=GREEN, gap=13)
hero_tile(s, "hero_solution", accent=GREEN)
page_no(s, 3)

# === Slide 4 — How it works (pipeline) =======================================
s = new_slide(prs)
header(s, "How it works", BLUE, ["A multi-agent pipeline on Google ADK"],
       hsize=33, hwidth=11.5)
stages = [
    ("Collect", "5 sources, parallel", "i_collect", BLUE),
    ("Normalize", "+ de-duplicate", "i_norm", RED),
    ("Process", "Gemini enrich", "i_process", YELLOW),
    ("Guardrail", "fact-check", "i_guard", GREEN),
    ("Digest", "EN + AR narrative", "i_digest", BLUE),
    ("Render", "xlsx · html · md", "i_render", RED),
]
n = len(stages)
gap = 0.26
usable = SW - 2 * 0.75
bw = (usable - (n - 1) * gap) / n
bh = 1.95
by = 3.05
bx0 = 0.75
for i, (name, sub, ic, accent) in enumerate(stages):
    bx = bx0 + i * (bw + gap)
    rect(s, bx, by, bw, bh, WHITE, line=LINE, line_w=1.25, rounded=True, radius=0.10)
    ih = 0.62
    pic(s, icon(ic), bx + (bw - ih) / 2, by + 0.26, h=ih)
    text(s, bx, by + 1.02, bw, 0.32, name, size=13.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, bx, by + 1.40, bw, 0.42, sub, size=9.5, color=MUTE,
         align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < n - 1:
        text(s, bx + bw, by + bh / 2 - 0.18, gap, 0.36, "›", size=22, color=FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
text(s, 0.75, 5.45, usable, 0.4,
     "Everything runs through one ADK agent tree — collectors fan out in parallel, "
     "then one Gemini pipeline produces a trustworthy digest.",
     size=14, color=MUTE, align=PP_ALIGN.CENTER)
page_no(s, 4)

# === Slide 5 — The agent tree (orchestration) ================================
s = new_slide(prs)
header(s, "Orchestration", BLUE, ["One ADK agent tree"], hsize=34, hwidth=11.5)
# left: vertical timeline of the SequentialAgent's stages
rail_x = 1.22
root_y = 2.52
nodes_y = [3.06, 3.58, 4.42, 4.92, 5.42, 5.92, 6.42]
rect(s, rail_x - 0.013, root_y, 0.026, nodes_y[-1] - root_y, RGBColor(0xBD, 0xC1, 0xC6))
oval(s, rail_x - 0.18, root_y - 0.18, 0.36, BLUE)
text(s, 1.58, root_y - 0.18, 3.1, 0.36, "NewsCatchUpPipeline", size=15.5, color=INK,
     bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, 4.80, root_y - 0.17, 3.0, 0.34, "SequentialAgent", size=12, color=BLUE,
     bold=True, anchor=MSO_ANCHOR.MIDDLE)
tree_nodes = [
    (1, "PipelineInit", "create & ID the run", FAINT),
    (2, "CollectSources", "ParallelAgent — fan-out", BLUE),
    (3, "NormalizeDedup", "merge · de-duplicate", RED),
    (4, "Processing", "Gemini enrich", AMBER),
    (5, "GuardrailCritic", "fact-check · self-correct", GREEN),
    (6, "DigestEditor", "EN / AR narrative", BLUE),
    (7, "Render", "xlsx · html · md", RED),
]
for (num, name, role, color), ny in zip(tree_nodes, nodes_y):
    tnode(s, rail_x, ny, num, name, role, color)
text(s, 1.58, 3.88, 6.4, 0.3,
     "RSS · Scrape · API · Search · YouTube — run concurrently",
     size=11.5, color=MUTE, anchor=MSO_ANCHOR.MIDDLE)
# right: count panel
px, py, pw, ph = 8.9, 2.52, 3.5, 4.04
rect(s, px, py, pw, ph, TILE, line=LINE, line_w=1.0, rounded=True, radius=0.06)
text(s, px + 0.36, py + 0.26, pw - 0.7, 0.95, "13", size=58, color=BLUE, bold=True)
text(s, px + 0.40, py + 1.28, pw - 0.7, 0.4, "agents · one tree", size=13.5, color=MUTE)
rect(s, px + 0.38, py + 1.82, pw - 0.76, 0.012, LINE)
tree_stats = [
    ("1", "Sequential orchestrator", BLUE),
    ("1", "Parallel fan-out", GREEN),
    ("5", "Source collectors", AMBER),
    ("6", "Pipeline stages", RED),
]
sy = py + 2.04
for val, label, c in tree_stats:
    text(s, px + 0.38, sy, 0.5, 0.32, val, size=16, color=c, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, px + 0.86, sy, pw - 1.2, 0.32, label, size=12, color=BODY,
         anchor=MSO_ANCHOR.MIDDLE)
    sy += 0.48
page_no(s)

# === Slide 6 — What each agent does ==========================================
s = new_slide(prs)
header(s, "The agents", BLUE, ["What each agent does"], hsize=34, hwidth=11.5)
lx, rx, col_w = 0.95, 7.05, 5.6
text(s, lx, 2.48, col_w, 0.3, "PIPELINE STAGES", size=12, color=BLUE, bold=True,
     spacing=1.8)
left_rows = [
    ("PipelineInit", "create & ID the run"),
    ("CollectSources", "5 collectors, in parallel"),
    ("NormalizeDedup", "merge & de-duplicate"),
    ("Processing", "enrich with Gemini"),
    ("GuardrailCritic", "fact-check & self-correct"),
    ("DigestEditor", "write the EN / AR story"),
    ("Render", "Excel · HTML · Markdown"),
]
ly = 3.0
for name, role in left_rows:
    kv_row(s, lx, ly, col_w, name, role, INK)
    ly += 0.52
text(s, rx, 2.48, col_w, 0.3, "GEMINI SUB-AGENTS  (LLM)", size=12, color=AMBER,
     bold=True, spacing=1.8)
right_rows = [
    ("news_processor", "categorize · score · summarize"),
    ("faithfulness_critic", "verify vs the source"),
    ("digest_editor", "compose the narrative"),
    ("search_collector", "Google Search grounding"),
    ("youtube_summarizer", "summarize transcripts"),
]
ry = 3.0
for name, role in right_rows:
    kv_row(s, rx, ry, col_w, name, role, AMBER)
    ry += 0.52
text(s, rx, ry + 0.10, col_w, 0.5,
     "+ news_reprocessor (self-correction)  ·  enrichment_judge (offline eval)",
     size=11, color=FAINT, line_spacing=1.15)
page_no(s)

# === Slide 7 — Architecture decisions ========================================
s = new_slide(prs)
header(s, "Architecture", BLUE, ["Four decisions behind the build"],
       hsize=33, hwidth=11.5)
cards = [
    ("ADK-native", "The agent tree IS the orchestration — no glue code.", BLUE),
    ("Free → Cloud by config", "Same code runs local or on Google Cloud, not a rewrite.", GREEN),
    ("Swappable storage", "SQLite ↔ Firestore behind a clean storage port; AI Studio ↔ Vertex by config.", AMBER),
    ("Secure by default", "Loopback-guarded settings, fail-closed auth & guardrails, one SSRF chokepoint.", RED),
]
gx, gy = ML, 2.6
cw, ch = (SW - 2 * ML - 0.45) / 2, 1.78
hgap, vgap = 0.45, 0.32
for i, (title, why, accent) in enumerate(cards):
    cx = gx + (i % 2) * (cw + hgap)
    cy = gy + (i // 2) * (ch + vgap)
    rect(s, cx, cy, cw, ch, TILE, line=LINE, line_w=1.0, rounded=True, radius=0.06)
    rect(s, cx, cy, 0.10, ch, accent, rounded=True, radius=0.5)
    text(s, cx + 0.40, cy + 0.30, cw - 0.7, 0.4, title, size=18.5, color=INK, bold=True)
    text(s, cx + 0.40, cy + 0.86, cw - 0.7, 0.8, why, size=13.5, color=MUTE,
         line_spacing=1.1)
page_no(s, 5)

# === Slide 6 — Feature 01: Collect ===========================================
s = new_slide(prs)
header(s, "Feature 01", BLUE, ["Collect from everywhere"], hsize=34, hwidth=11.5)
text(s, ML, 2.48, 11.3, 0.7,
     "One safe pipeline pulls from five source types — SSRF-guarded and rate-limited.",
     size=16.5, color=MUTE, line_spacing=1.1)
sources = [
    (icon("src_rss"), "RSS"),
    (icon("src_web"), "Web scrape"),
    (icon("src_api"), "News APIs"),
    (icon("src_youtube"), "YouTube"),
    (icon("src_search"), "Search grounding"),
]
# tiles row
n = len(sources)
tw = 1.9
tgap = 0.42
total = n * tw + (n - 1) * tgap
sx = SW / 2 - total / 2
ty = 3.55
for i, (p, label) in enumerate(sources):
    x = sx + i * (tw + tgap)
    rect(s, x, ty, tw, 1.75, TILE, line=LINE, line_w=1.0, rounded=True, radius=0.10)
    ih = 0.78
    pic(s, p, x + (tw - ih) / 2, ty + 0.34, h=ih)
    text(s, x, ty + 1.24, tw, 0.4, label, size=12.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
why_note(s, BLUE, "No single source has the whole picture.")
page_no(s, 6)

# === Slide 7 — Feature 02: Rank & summarize ==================================
s = new_slide(prs)
header(s, "Feature 02", AMBER, ["Rank & summarize", "with Gemini"])
hero_bullets(s, [
    "Categorize into AI/Tech, Business, World, and Gulf/MENA",
    "Score importance — see what matters first",
    "Summarize in English + Arabic",
    "Extract the people, orgs, and places",
], size=16.5, color=AMBER, gap=11)
hero_tile(s, "i_process", accent=AMBER)
why_note(s, AMBER, "Triage-first — the signal, in your language.")
page_no(s, 7)

# === Slide 8 — Feature 03: Guardrail =========================================
s = new_slide(prs)
header(s, "Feature 03", GREEN, ["An agent you can trust"])
hero_bullets(s, [
    "An LLM-as-judge critic fact-checks top items against their source",
    "Unfaithful summaries are flagged — and never shown",
    "The judge and the critic share one faithfulness rubric",
], size=17, color=GREEN, gap=15)
hero_tile(s, "i_guard", accent=GREEN)
why_note(s, GREEN, "No hallucinations. No made-up news.")
page_no(s, 8)

# === Slide 9 — Feature 04: Delivery ==========================================
s = new_slide(prs)
header(s, "Feature 04", BLUE, ["Delivered the way you work"], hsize=34)
hero_bullets(s, [
    "Excel · HTML dashboard · Markdown digests",
    "A Next.js console to browse, search, and configure",
    "A one-click macOS desktop app",
], size=17, color=BLUE, gap=15)
hero_tile(s, "hero_delivery", accent=BLUE)
why_note(s, BLUE, "Meet every user on the surface they already use.")
page_no(s, 9)

# === Slide 12 — How I built it (process) =====================================
s = new_slide(prs)
header(s, "Process", BLUE, ["How I directed the build"], hsize=34, hwidth=11.5)
# top: the build loop
steps = [
    ("Brainstorm", "intent & design", "p_brainstorm", BLUE),
    ("Spec", "write it down", "p_spec", RED),
    ("Plan", "Codex-reviewed", "p_plan", AMBER),
    ("Build", "TDD · subagents", "p_build", GREEN),
    ("Review", "Codex gate", "p_review", BLUE),
    ("Merge", "green & verified", "p_merge", GREEN),
]
n = len(steps)
gap = 0.24
usable = SW - 2 * 0.75
bw = (usable - (n - 1) * gap) / n
bh = 1.5
by = 2.46
for i, (name, sub, ic, accent) in enumerate(steps):
    bx = 0.75 + i * (bw + gap)
    rect(s, bx, by, bw, bh, WHITE, line=LINE, line_w=1.25, rounded=True, radius=0.10)
    ih = 0.5
    pic(s, icon(ic), bx + (bw - ih) / 2, by + 0.22, h=ih)
    text(s, bx, by + 0.80, bw, 0.3, name, size=12.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, bx, by + 1.13, bw, 0.3, sub, size=9, color=MUTE, align=PP_ALIGN.CENTER)
    if i < n - 1:
        text(s, bx + bw, by + bh / 2 - 0.18, gap, 0.36, "›", size=20, color=FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
text(s, 0.75, by + bh + 0.14, usable, 0.3,
     "I set the architecture & standards up front, then looped every feature "
     "through this — reviewed and tested before it merged.",
     size=12.5, color=MUTE, align=PP_ALIGN.CENTER)
# bottom: the 3 things that made it work
princ = [
    ("Standards up front",
     "Architecture, security & the 'Signal' design system — fixed in CLAUDE.md before any code.", BLUE),
    ("Independent review gates",
     "Every plan and task passed an external Codex review before merge.", GREEN),
    ("Persistent memory",
     "A living BUILD-LOG and memory kept the agent on-thread across sessions.", AMBER),
]
cy = 4.86
cw = (SW - 2 * ML - 2 * 0.4) / 3
ch = 1.5
for i, (title, body, accent) in enumerate(princ):
    cx = ML + i * (cw + 0.4)
    rect(s, cx, cy, cw, ch, TILE, line=LINE, line_w=1.0, rounded=True, radius=0.07)
    rect(s, cx, cy, 0.09, ch, accent, rounded=True, radius=0.5)
    text(s, cx + 0.34, cy + 0.24, cw - 0.6, 0.4, title, size=14.5, color=INK, bold=True)
    text(s, cx + 0.34, cy + 0.72, cw - 0.58, 0.7, body, size=11.5, color=MUTE,
         line_spacing=1.12)
page_no(s)

# === Slide 10 — Closing ======================================================
s = new_slide(prs)
header(s, "Summary", GREEN, ["Free today.", "Cloud-ready tomorrow."])
bullets(s, ML, 3.55, 7.4, 2.0, [
    "Runs free on your laptop",
    "Scales to Cloud Run + Vertex AI + Firestore — by config",
    "Built on Google ADK + Gemini",
], size=17, dash_color=GREEN, gap=12)
text(s, SW - ML - 4.6, 1.55, 4.6, 0.3, "BUILT WITH", size=11, color=FAINT,
     bold=True, spacing=2.6, align=PP_ALIGN.LEFT)
logo_strip(s, TECH, 2.05, icon_h=0.62, cell_w=0.78, center_x=SW - ML - 2.3,
           labels=True, label_size=9)
rect(s, ML, 6.5, SW - 2 * ML, 0.013, LINE)
text(s, ML, 6.66, 8, 0.3, "Catch-Up — News Intelligence Agent", size=12, color=MUTE, bold=True)
text(s, SW - ML - 5, 6.66, 5, 0.3, "Ahmed Hesham · AI Engineer", size=12, color=MUTE,
     align=PP_ALIGN.RIGHT)
page_no(s, 10)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
